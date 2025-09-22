#!/usr/bin/env python3
"""
PPT 파일에서 슬라이드 이미지를 추출하는 스크립트
"""

import os
from pptx import Presentation
from PIL import Image
import io

def extract_slides_from_ppt(ppt_path, output_dir="slides"):
    """PPT 파일에서 각 슬라이드를 이미지로 추출"""
    
    # 출력 디렉토리 생성
    os.makedirs(output_dir, exist_ok=True)
    
    try:
        prs = Presentation(ppt_path)
        print(f"PPT 파일 로딩 완료: {ppt_path}")
        print(f"총 {len(prs.slides)}개 슬라이드 발견")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"슬라이드 {slide_num} 처리 중...")
            
            # 슬라이드를 이미지로 변환
            slide_image = slide_to_image(slide)
            
            # 이미지 저장
            output_path = os.path.join(output_dir, f"slide_{slide_num}.png")
            slide_image.save(output_path, "PNG")
            print(f"저장 완료: {output_path}")
        
        print(f"\n모든 슬라이드 추출 완료! 결과: {output_dir}/")
        return True
        
    except Exception as e:
        print(f"PPT 처리 중 오류 발생: {e}")
        return False

def slide_to_image(slide):
    """슬라이드를 PIL Image로 변환"""
    # 슬라이드의 크기 정보 가져오기
    slide_width = slide.slide_width
    slide_height = slide.slide_height
    
    # 빈 이미지 생성 (흰색 배경)
    img = Image.new('RGB', (slide_width, slide_height), 'white')
    
    # 여기서는 간단한 텍스트만 추출하여 이미지에 그리기
    # 실제로는 더 복잡한 슬라이드 내용을 처리해야 할 수 있습니다
    from PIL import ImageDraw, ImageFont
    
    draw = ImageDraw.Draw(img)
    
    # 기본 폰트 사용
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
    except:
        font = ImageFont.load_default()
    
    y_position = 50
    
    # 슬라이드의 텍스트 추출 및 그리기
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()
            # 텍스트를 여러 줄로 나누기
            lines = text.split('\n')
            for line in lines:
                if line.strip():
                    draw.text((50, y_position), line.strip(), fill='black', font=font)
                    y_position += 30
    
    return img

if __name__ == "__main__":
    import sys
    
    # 명령행 인수로 PPT 파일 경로 받기
    if len(sys.argv) > 1:
        ppt_file = sys.argv[1]
    else:
        # 기본값으로 현재 디렉토리에서 PPT 파일 찾기
        ppt_files = [f for f in os.listdir('.') if f.endswith(('.pptx', '.ppt'))]
        if ppt_files:
            ppt_file = ppt_files[0]
        else:
            print("❌ PPT 파일을 찾을 수 없습니다. .pptx 또는 .ppt 파일이 필요합니다.")
            sys.exit(1)
    
    if os.path.exists(ppt_file):
        print(f"📊 PPT 파일 처리 시작: {ppt_file}")
        extract_slides_from_ppt(ppt_file)
    else:
        print(f"❌ PPT 파일을 찾을 수 없습니다: {ppt_file}")
        sys.exit(1)
