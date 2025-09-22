import os
import fitz  # PyMuPDF
import base64
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches
import io
from PIL import Image
import torch
import torchaudio
import soundfile as sf
import numpy as np
from zonos.model import Zonos
from zonos.conditioning import make_cond_dict
from zonos.utils import DEFAULT_DEVICE as device

# 1. API 키를 환경 변수에서 안전하게 불러오기
api_key = os.environ.get("OPENAI_API_KEY")
if not api_key:
    raise ValueError("OPENAI_API_KEY 환경 변수가 설정되지 않았습니다. API 키를 설정해주세요.")
client = OpenAI(api_key=api_key)

# Zonos 모델 초기화 (전역 변수로 관리)
ZONOS_MODEL = None
SPEAKER_EMBEDDING = None


def initialize_zonos_model(model_name="Zyphra/Zonos-v0.1-transformer"):
    """Zonos 모델을 초기화하는 함수"""
    global ZONOS_MODEL
    if ZONOS_MODEL is None:
        print("Zonos 모델을 로딩 중...")
        ZONOS_MODEL = Zonos.from_pretrained(model_name, device=device)
        ZONOS_MODEL.requires_grad_(False).eval()
        print("Zonos 모델 로딩 완료!")
    return ZONOS_MODEL


def load_speaker_embedding(speaker_audio_path):
    """샘플 오디오에서 스피커 임베딩을 생성하는 함수"""
    global SPEAKER_EMBEDDING
    if SPEAKER_EMBEDDING is None and speaker_audio_path:
        print(f"스피커 오디오에서 임베딩 생성 중: {speaker_audio_path}")
        try:
            # torchaudio를 사용하여 오디오 로드 (Gradio 방식과 동일)
            wav, sampling_rate = torchaudio.load(speaker_audio_path)
            
            # 스피커 임베딩 생성
            SPEAKER_EMBEDDING = ZONOS_MODEL.make_speaker_embedding(wav, sampling_rate)
            # bfloat16으로 변환 (Gradio 방식과 동일)
            SPEAKER_EMBEDDING = SPEAKER_EMBEDDING.to(device, dtype=torch.bfloat16)
            print("스피커 임베딩 생성 완료!")
        except Exception as e:
            print(f"torchaudio 로딩 실패, soundfile로 시도: {e}")
            try:
                # soundfile을 사용하여 오디오 로드
                wav, sampling_rate = sf.read(speaker_audio_path)
                # numpy 배열을 torch tensor로 변환
                wav = torch.from_numpy(wav).float()
                # 모노 채널로 변환 (필요한 경우)
                if wav.dim() > 1:
                    wav = wav.mean(dim=-1)
                # 배치 차원 추가
                wav = wav.unsqueeze(0)
                
                SPEAKER_EMBEDDING = ZONOS_MODEL.make_speaker_embedding(wav, sampling_rate)
                SPEAKER_EMBEDDING = SPEAKER_EMBEDDING.to(device, dtype=torch.bfloat16)
                print("스피커 임베딩 생성 완료! (soundfile 사용)")
            except Exception as e2:
                print(f"soundfile로도 로딩 실패: {e2}")
                return None
    return SPEAKER_EMBEDDING


def generate_voice_from_text(text, speaker_audio_path=None, language="en-us", output_path=None):
    """텍스트를 음성으로 변환하는 함수"""
    # 모델 초기화
    model = initialize_zonos_model()
    
    # 스피커 임베딩 로드
    if speaker_audio_path:
        speaker_embedding = load_speaker_embedding(speaker_audio_path)
    else:
        speaker_embedding = None
    
    # 텍스트 길이 확인 (디버깅용)
    print(f"처리할 전체 텍스트 길이: {len(text)}자")
    
    # 긴 텍스트의 경우 문장 단위로 분할하여 처리
    if len(text) > 800:
        print(f"긴 텍스트 감지 ({len(text)}자). 문장 단위로 분할하여 처리합니다.")
        return generate_long_text_voice(text, model, speaker_embedding, language, output_path)
    
    # Gradio 방식의 speaking_rate 계산
    text_length = len(text)
    
    # Gradio에서 사용하는 공식: estimated_generation_duration = 30 * len(text) / 400
    # 이를 역산하여 speaking_rate 계산
    estimated_duration = 30 * text_length / 400  # Gradio 방식
    estimated_duration = max(min(estimated_duration, 30), 5)  # 5-30초 범위로 제한
    
    # speaking_rate는 초당 음소 수이므로, 텍스트 길이를 기반으로 계산
    # 한국어의 경우 대략적으로 글자 수 * 1.5 = 음소 수
    estimated_phonemes = text_length * 1.5
    speaking_rate = estimated_phonemes / estimated_duration
    
    # speaking_rate 범위 제한 (5-30)
    speaking_rate = max(min(speaking_rate, 30), 5)
    
    print(f"텍스트 길이: {text_length}자, 예상 길이: {estimated_duration:.1f}초, Speaking rate: {speaking_rate:.1f}")
    
    # 조건부 생성 설정 - 보이스 클로닝 최적화
    cond_dict = make_cond_dict(
        text=text,
        speaker=speaker_embedding,
        language=language,
        device=device,
        # 보이스 클로닝 최적화 파라미터
        pitch_std=45.0,  # 자연스러운 음성 (Gradio 기본값)
        speaking_rate=speaking_rate,  # 텍스트 길이에 따라 조정
        vqscore_8=[0.78] * 8,  # Gradio 기본값
        dnsmos_ovrl=4.0,  # Gradio 기본값
        fmax=22050.0,  # 보이스 클로닝에 최적화된 주파수
        speaker_noised=False  # 스피커 노이즈 제거 비활성화
    )
    conditioning = model.prepare_conditioning(cond_dict)
    
    # 음성 생성
    print(f"음성 생성 중: '{text[:50]}...'")
    codes = model.generate(conditioning)
    wavs = model.autoencoder.decode(codes).cpu()
    
    # 오디오 저장
    if output_path:
        try:
            # soundfile을 사용하여 오디오 저장
            # 오디오 데이터를 numpy 배열로 변환하고 정규화
            audio_data = wavs[0].numpy()
            # 오디오 데이터가 2D인 경우 1D로 변환
            if audio_data.ndim > 1:
                audio_data = audio_data.squeeze()
            # 데이터 타입을 float32로 변환
            audio_data = audio_data.astype(np.float32)
            # 샘플링 레이트를 정수로 변환
            sample_rate = int(model.autoencoder.sampling_rate)
            
            sf.write(output_path, audio_data, sample_rate, format='WAV', subtype='PCM_16')
            print(f"오디오 저장 완료: {output_path}")
        except Exception as e:
            print(f"soundfile 저장 실패, torchaudio로 시도: {e}")
            try:
                torchaudio.save(output_path, wavs[0], model.autoencoder.sampling_rate)
                print(f"오디오 저장 완료: {output_path} (torchaudio 사용)")
            except Exception as e2:
                print(f"오디오 저장 실패: {e2}")
    
    return wavs[0], model.autoencoder.sampling_rate


def generate_long_text_voice(text, model, speaker_embedding, language, output_path):
    """긴 텍스트를 문장 단위로 분할하여 음성 생성하는 함수"""
    # 문장 단위로 분할 (마침표 기준)
    sentences = [s.strip() for s in text.split('.') if s.strip()]
    
    print(f"총 {len(sentences)}개 문장으로 분할")
    
    all_audio = []
    sample_rate = None
    
    for i, sentence in enumerate(sentences):
        if not sentence.endswith('.'):
            sentence += '.'
        
        print(f"[{i+1}/{len(sentences)}] 문장 처리: '{sentence[:30]}...'")
        
        # 각 문장에 대한 speaking_rate 계산
        text_length = len(sentence)
        estimated_phonemes = text_length * 1.5
        target_duration = min(max(estimated_phonemes / 15, 3), 15)  # 3-15초 범위
        speaking_rate = estimated_phonemes / target_duration
        speaking_rate = max(min(speaking_rate, 30), 5)
        
        # 조건부 생성 설정 - 보이스 클로닝 최적화
        cond_dict = make_cond_dict(
            text=sentence,
            speaker=speaker_embedding,
            language=language,
            device=device,
            pitch_std=20.0,  # 자연스러운 음성
            speaking_rate=speaking_rate,
            vqscore_8=[0.78] * 8,  # Gradio 기본값
            dnsmos_ovrl=4.0,  # Gradio 기본값
            fmax=22050.0,  # 보이스 클로닝에 최적화된 주파수
            speaker_noised=False  # 스피커 노이즈 제거 비활성화
        )
        conditioning = model.prepare_conditioning(cond_dict)
        
        # 음성 생성
        codes = model.generate(conditioning)
        wavs = model.autoencoder.decode(codes).cpu()
        
        all_audio.append(wavs[0])
        if sample_rate is None:
            sample_rate = model.autoencoder.sampling_rate
    
    # 모든 오디오를 연결
    if len(all_audio) > 1:
        combined_audio = torch.cat(all_audio, dim=-1)
    else:
        combined_audio = all_audio[0]
    
    # 오디오 저장
    if output_path:
        try:
            audio_data = combined_audio.numpy()
            if audio_data.ndim > 1:
                audio_data = audio_data.squeeze()
            audio_data = audio_data.astype(np.float32)
            sample_rate = int(sample_rate)
            
            sf.write(output_path, audio_data, sample_rate, format='WAV', subtype='PCM_16')
            print(f"긴 텍스트 오디오 저장 완료: {output_path}")
        except Exception as e:
            print(f"긴 텍스트 오디오 저장 실패: {e}")
    
    return combined_audio, sample_rate


def generate_voice_for_scripts(scripts_folder, speaker_audio_path=None, language="en-us"):
    """스크립트 폴더의 모든 텍스트 파일을 음성으로 변환하는 함수"""
    if not os.path.exists(scripts_folder):
        print(f"스크립트 폴더가 존재하지 않습니다: {scripts_folder}")
        return
    
    # 오디오 저장 폴더 생성
    audio_folder = f"{scripts_folder}_audio"
    os.makedirs(audio_folder, exist_ok=True)
    print(f"오디오가 '{audio_folder}' 폴더에 저장됩니다.")
    
    # 스크립트 파일들 찾기
    script_files = [f for f in os.listdir(scripts_folder) if f.endswith('.txt')]
    script_files.sort()  # 파일명 순으로 정렬
    
    if not script_files:
        print("스크립트 파일을 찾을 수 없습니다.")
        return
    
    print(f"총 {len(script_files)}개의 스크립트 파일을 처리합니다.")
    
    # 각 스크립트 파일 처리
    for i, script_file in enumerate(script_files, 1):
        script_path = os.path.join(scripts_folder, script_file)
        print(f"\n[{i}/{len(script_files)}] {script_file} 처리 중...")
        
        try:
            # 스크립트 파일 읽기
            with open(script_path, 'r', encoding='utf-8') as f:
                script_text = f.read().strip()
            
            if not script_text:
                print(f"빈 스크립트 파일: {script_file}")
                continue
            
            # 오디오 파일명 생성
            audio_filename = script_file.replace('.txt', '.wav')
            audio_path = os.path.join(audio_folder, audio_filename)
            
            # 메모리 정리 (매 5개 파일마다)
            if i % 5 == 0:
                print("메모리 정리 중...")
                import gc
                gc.collect()
                if torch.cuda.is_available():
                    torch.cuda.empty_cache()
            
            # 음성 생성
            generate_voice_from_text(
                text=script_text,
                speaker_audio_path=speaker_audio_path,
                language=language,
                output_path=audio_path
            )
            
        except Exception as e:
            print(f"❌ {script_file} 처리 중 오류 발생: {e}")
    
    print(f"\n🎉 모든 스크립트의 음성 생성이 완료되었습니다. 결과: {audio_folder}")


def extract_text_and_images_from_pdf(pdf_path):
    """PDF 파일에서 각 페이지의 텍스트와 이미지를 추출하는 함수"""
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(pdf_path)
        pages_data = []
        
        for page_num in range(doc.page_count):
            page = doc[page_num]
            page_data = {
                'page_number': page_num + 1,
                'text': '',
                'images': []
            }
            
            # 텍스트 추출
            text = page.get_text()
            page_data['text'] = text
            
            # 페이지를 이미지로 변환하여 Base64로 인코딩
            try:
                # 고해상도 이미지로 변환
                zoom = 2  # 200% 확대 (200 DPI)
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)
                
                # 이미지를 bytes로 변환
                img_bytes = pix.tobytes("png")
                
                # Base64로 인코딩
                base64_image = base64.b64encode(img_bytes).decode('utf-8')
                page_data['images'].append(base64_image)
                
            except Exception as e:
                print(f"페이지 {page_num + 1}의 이미지 처리 중 오류: {e}")
                continue
            
            pages_data.append(page_data)
        
        doc.close()
        return pages_data
    
    except Exception as e:
        print(f"PDF 파일 처리 중 오류 발생: {e}")
        return []


def generate_presentation_script_for_pdf(pdf_path, prompt):
    """PDF의 모든 페이지를 처리하여 발표 대본을 생성하고 지정된 폴더에 파일로 저장하는 함수"""
    
    # PDF 파일명에서 확장자를 제거하여 기본 이름 추출
    pdf_base_name = os.path.splitext(os.path.basename(pdf_path))[0]
    
    # 결과를 저장할 폴더 이름 지정
    output_folder = f"{pdf_base_name}_scripts"
    
    # 폴더가 존재하지 않으면 생성
    os.makedirs(output_folder, exist_ok=True)
    print(f"결과가 '{output_folder}' 폴더에 저장됩니다.")

    # PDF에서 페이지 데이터 추출
    pages_data = extract_text_and_images_from_pdf(pdf_path)
    
    if not pages_data:
        print("PDF 파일에서 데이터를 추출할 수 없습니다.")
        return
    
    total_pages = len(pages_data)
    print(f"총 {total_pages} 페이지의 PDF 파일을 처리합니다.")

    # 각 페이지를 순회하며 대본 생성
    for page_data in pages_data:
        page_num = page_data['page_number']
        print(f"\n[{page_num}/{total_pages}] 페이지 처리 중...")
        
        text = page_data['text']
        images = page_data['images']

        # 간단한 프롬프트 (두 문장 내로)
        page_prompt = f"이 페이지의 내용을 두 문장 내로 간결하게 설명하는 발표 대본을 작성해주세요. {prompt}"

        # GPT-4o API에 텍스트와 이미지 함께 전송
        try:
            content = [
                {"type": "text", "text": f"{page_prompt}\n\n---페이지 내용---\n{text}"}
            ]
            
            # 이미지가 있으면 추가
            for img_base64 in images:
                content.append({
                    "type": "image_url",
                    "image_url": {
                        "url": f"data:image/png;base64,{img_base64}"
                    }
                })
            
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": content
                    }
                ],
                max_tokens=200,  # 짧은 대본을 위해 토큰 수 제한
            )
            
            script_content = response.choices[0].message.content
            
            # 생성된 대본을 지정된 폴더 안에 파일로 저장
            output_filename = f"{pdf_base_name}_page_{page_num}_script.txt"
            output_path = os.path.join(output_folder, output_filename)
            
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(script_content)
            
            print(f"✅ '{output_path}' 파일 저장 완료.")

        except Exception as e:
            print(f"❌ {page_num} 페이지 처리 중 API 오류 발생: {e}")

    print("\n🎉 모든 페이지의 발표 대본 생성이 완료되었습니다.")


# def generate_presentation_script_for_pdf(pdf_path, prompt):
#     """PDF의 모든 페이지를 처리하여 발표 대본을 생성하고 지정된 폴더에 파일로 저장하는 함수"""
#     
#     # PDF 파일명에서 확장자를 제거하여 기본 이름 추출 (예: "중간보고서 미팅")
#     pdf_base_name = os.path.splitext(os.path.basename(pdf_path))[0]
#     
#     # 결과를 저장할 폴더 이름 지정 (예: "중간보고서 미팅_scripts")
#     output_folder = f"{pdf_base_name}_scripts"
#     
#     # 폴더가 존재하지 않으면 생성
#     os.makedirs(output_folder, exist_ok=True)
#     print(f"결과가 '{output_folder}' 폴더에 저장됩니다.")

#     try:
#         doc = fitz.open(pdf_path)
#         total_pages = doc.page_count
#         print(f"총 {total_pages} 페이지의 PDF 파일을 처리합니다.")

#     except Exception as e:
#         print(f"PDF 파일을 여는 중 오류 발생: {e}")
#         return

#     # 각 페이지를 순회하며 대본 생성
#     for page_num in range(total_pages):
#         current_page = page_num + 1
#         print(f"\n[{current_page}/{total_pages}] 페이지 처리 중...")
#         
#         page = doc[page_num]

#         # 1. 현재 페이지의 텍스트 추출
#         text = page.get_text()

#         # 2. 현재 페이지를 이미지로 변환 및 Base64 인코딩
#         pix = page.get_pixmap()
#         image_bytes = pix.tobytes("png")
#         base64_image = base64.b64encode(image_bytes).decode('utf-8')

#         # 3. GPT-4o API에 텍스트와 이미지 함께 전송
#         try:
#             response = client.chat.completions.create(
#                 model="gpt-4o",
#                 messages=[
#                     {
#                         "role": "user",
#                         "content": [
#                             {"type": "text", "text": f"이것은 전체 발표 자료 중 {current_page}번째 페이지입니다. {prompt}\n\n---추출된 텍스트---\n{text}"},
#                             {
#                                 "type": "image_url",
#                                 "image_url": {
#                                     "url": f"data:image/png;base64,{base64_image}"
#                                 }
#                             }
#                         ]
#                     }
#                 ],
#                 max_tokens=1000,
#             )
#             
#             script_content = response.choices[0].message.content
            
#             # 4. 생성된 대본을 지정된 폴더 안에 파일로 저장
#             output_filename = f"{pdf_base_name}_page_{current_page}_script.txt"
#             output_path = os.path.join(output_folder, output_filename)
#             
#             with open(output_path, "w", encoding="utf-8") as f:
#                 f.write(script_content)
#             
#             print(f"✅ '{output_path}' 파일 저장 완료.")

#         except Exception as e:
#             print(f"❌ {current_page} 페이지 처리 중 API 오류 발생: {e}")

#     doc.close()
#     print("\n🎉 모든 페이지의 발표 대본 생성이 완료되었습니다.")


def generate_presentation_script(file_path, prompt, speaker_audio_path=None, language="ko", generate_audio=True):
    """파일 확장자에 따라 PDF 또는 PPT를 자동으로 처리하고 선택적으로 음성을 생성하는 통합 함수"""
    
    # 파일 확장자 확인
    file_extension = os.path.splitext(file_path)[1].lower()
    
    if file_extension == '.pdf':
        print("📄 PDF 파일을 감지했습니다. PDF 처리 모드로 실행합니다.")
        generate_presentation_script_for_pdf(file_path, prompt)
    else:
        print(f"❌ 지원하지 않는 파일 형식입니다: {file_extension}")
        print("지원되는 형식: .pdf")
        return
    
    # 스크립트 생성 후 음성 생성
    if generate_audio:
        # 스크립트 폴더 경로 생성
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        scripts_folder = f"{base_name}_scripts"
        
        if os.path.exists(scripts_folder):
            print(f"\n🎵 스크립트에서 음성 생성 시작...")
            if speaker_audio_path and os.path.exists(speaker_audio_path):
                print(f"스피커 오디오 사용: {speaker_audio_path}")
            else:
                print("스피커 오디오가 제공되지 않았습니다. 기본 음성으로 생성합니다.")
            
            generate_voice_for_scripts(scripts_folder, speaker_audio_path, language)
        else:
            print(f"스크립트 폴더를 찾을 수 없습니다: {scripts_folder}")




def extract_pages_from_pdf(pdf_path, output_dir="slides"):
    """PDF 파일에서 각 페이지를 이미지로 추출하는 함수 (test_img.py 기반)"""
    try:
        import fitz  # PyMuPDF
        
        # 출력 디렉토리 생성
        os.makedirs(output_dir, exist_ok=True)
        
        print(f"📊 PDF 파일 처리 시작: {pdf_path}")
        
        # PDF 파일 열기
        pdf_document = fitz.open(pdf_path)
        
        print(f"총 {pdf_document.page_count} 페이지를 이미지로 변환합니다.")
        
        page_images = []
        
        for page_num in range(pdf_document.page_count):
            # 페이지 로드
            page = pdf_document.load_page(page_num)
            
            # 렌더링을 위한 설정 (dpi를 높여 고해상도 이미지로 만듭니다)
            zoom = 2  # 200% 확대 (200 DPI)
            mat = fitz.Matrix(zoom, zoom)
            
            # 페이지를 픽셀맵으로 변환 (이미지 데이터)
            pix = page.get_pixmap(matrix=mat)
            
            # 이미지 파일 경로 설정
            image_path = os.path.join(output_dir, f"slide_{page_num + 1}.png")
            
            # 픽셀맵을 PNG 파일로 저장
            pix.save(image_path)
            
            page_images.append(image_path)
            print(f"페이지 {page_num + 1}이 {image_path}에 저장되었습니다.")
        
        pdf_document.close()
        print("✅ 모든 페이지 변환이 완료되었습니다.")
        
        return page_images
        
    except Exception as e:
        print(f"❌ PDF 처리 중 오류 발생: {e}")
        return []




def convert_ppt_to_images_alternative(ppt_path, output_dir):
    """LibreOffice가 없을 때 사용하는 대체 방법"""
    try:
        print("🔄 대체 방법으로 슬라이드 이미지 생성 중...")
        
        from pptx import Presentation
        from PIL import Image, ImageDraw, ImageFont
        
        prs = Presentation(ppt_path)
        print(f"총 {len(prs.slides)}개 슬라이드 발견")
        
        slide_images = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"슬라이드 {slide_num} 처리 중...")
            
            # 슬라이드를 이미지로 변환 (텍스트 기반이지만 더 나은 레이아웃)
            slide_image = create_enhanced_slide_image(slide)
            
            if slide_image:
                # 이미지 저장
                output_path = os.path.join(output_dir, f"slide_{slide_num}.png")
                slide_image.save(output_path, "PNG")
                slide_images.append(output_path)
                print(f"저장 완료: {output_path}")
            else:
                print(f"❌ 슬라이드 {slide_num} 이미지 변환 실패")
        
        return slide_images
        
    except Exception as e:
        print(f"❌ 대체 방법 변환 중 오류: {e}")
        return []


def slide_to_actual_image(slide):
    """슬라이드를 실제 이미지로 변환하는 함수 (LibreOffice 사용)"""
    try:
        import subprocess
        import tempfile
        import os
        from PIL import Image
        
        # 임시 파일 생성
        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
            temp_path = tmp_file.name
        
        # LibreOffice를 사용해서 슬라이드를 이미지로 변환
        # 이 방법은 실제 슬라이드의 모든 요소를 포함한 이미지를 생성
        try:
            # LibreOffice 명령어로 슬라이드를 이미지로 변환
            cmd = [
                'libreoffice', '--headless', '--convert-to', 'png',
                '--outdir', os.path.dirname(temp_path),
                '--', temp_path
            ]
            
            # 이 방법은 복잡하므로, 더 간단한 방법 사용
            # 실제로는 python-pptx로는 완전한 슬라이드 이미지 추출이 제한적
            # 대신 텍스트와 기본 레이아웃을 포함한 이미지 생성
            
            return create_slide_image_with_layout(slide)
            
        except Exception as e:
            print(f"LibreOffice 변환 실패, 대체 방법 사용: {e}")
            return create_slide_image_with_layout(slide)
            
    except Exception as e:
        print(f"슬라이드 이미지 변환 중 오류: {e}")
        return None


def create_slide_image_with_layout(slide):
    """슬라이드의 레이아웃을 고려한 이미지 생성 (테스트 완료된 버전)"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        # 표준 슬라이드 크기 (16:9)
        slide_width = 1920
        slide_height = 1080
        
        # 이미지 생성 (흰색 배경)
        img = Image.new('RGB', (slide_width, slide_height), 'white')
        draw = ImageDraw.Draw(img)
        
        # 폰트 설정
        try:
            title_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 48)
            body_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 32)
        except:
            title_font = ImageFont.load_default()
            body_font = ImageFont.load_default()
        
        y_position = 120
        margin = 120
        
        # 슬라이드의 모든 텍스트를 레이아웃에 맞게 그리기
        for i, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                
                # 첫 번째 텍스트는 제목으로 처리
                if i == 0:
                    font = title_font
                    color = 'darkblue'
                    y_spacing = 60
                else:
                    font = body_font
                    color = 'black'
                    y_spacing = 40
                
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        # 텍스트가 화면을 넘지 않도록 처리
                        if len(line) > 50:
                            words = line.split()
                            current_line = ""
                            for word in words:
                                if len(current_line + " " + word) <= 50:
                                    current_line += " " + word if current_line else word
                                else:
                                    if current_line:
                                        draw.text((margin, y_position), current_line.strip(), fill=color, font=font)
                                        y_position += y_spacing
                                    current_line = word
                            if current_line:
                                draw.text((margin, y_position), current_line.strip(), fill=color, font=font)
                                y_position += y_spacing
                        else:
                            draw.text((margin, y_position), line.strip(), fill=color, font=font)
                            y_position += y_spacing
                        
                        # 슬라이드 높이를 넘지 않도록 제한
                        if y_position > slide_height - 100:
                            break
        
        return img
        
    except Exception as e:
        print(f"슬라이드 레이아웃 이미지 생성 중 오류: {e}")
        return None


def slide_to_image_python(slide, slide_width=1920, slide_height=1080):
    """슬라이드를 PIL Image로 변환하는 함수 (간단한 텍스트 기반)"""
    from PIL import Image, ImageDraw, ImageFont
    
    # 안전한 크기로 제한
    slide_width = min(slide_width, 1920)
    slide_height = min(slide_height, 1080)
    
    # 빈 이미지 생성 (흰색 배경)
    img = Image.new('RGB', (slide_width, slide_height), 'white')
    draw = ImageDraw.Draw(img)
    
    # 기본 폰트 사용 (크기 조정)
    font_size = max(24, min(48, slide_width // 50))  # 화면 크기에 맞게 폰트 크기 조정
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
    except:
        try:
            font = ImageFont.truetype("arial.ttf", font_size)
        except:
            font = ImageFont.load_default()
    
    y_position = 80
    margin = 80
    line_height = font_size + 10
    
    # 슬라이드의 텍스트 추출 및 그리기
    all_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            all_text.append(shape.text.strip())
    
    # 모든 텍스트를 합쳐서 처리
    combined_text = "\n".join(all_text)
    lines = combined_text.split('\n')
    
    for line in lines:
        if line.strip():
            # 텍스트가 이미지 너비를 넘지 않도록 처리
            max_chars = slide_width // (font_size // 2)  # 대략적인 문자 수 계산
            
            if len(line) > max_chars:
                # 긴 텍스트는 여러 줄로 나누기
                words = line.split()
                current_line = ""
                for word in words:
                    if len(current_line + " " + word) <= max_chars:
                        current_line += " " + word if current_line else word
                    else:
                        if current_line:
                            draw.text((margin, y_position), current_line.strip(), fill='black', font=font)
                            y_position += line_height
                        current_line = word
                if current_line:
                    draw.text((margin, y_position), current_line.strip(), fill='black', font=font)
                    y_position += line_height
            else:
                draw.text((margin, y_position), line.strip(), fill='black', font=font)
                y_position += line_height
            
            # 슬라이드 높이를 넘지 않도록 제한
            if y_position > slide_height - 80:
                break
    
    return img


def get_audio_duration(audio_path):
    """오디오 파일의 길이를 가져오는 함수"""
    import subprocess
    
    try:
        result = subprocess.run([
            'ffprobe', '-v', 'error', '-show_entries', 'format=duration',
            '-of', 'default=noprint_wrappers=1', audio_path
        ], capture_output=True, text=True)
        
        if result.returncode == 0:
            duration = float(result.stdout.strip())
            return int(duration)  # 정수로 반환
        else:
            print(f"❌ 오디오 길이 가져오기 실패: {audio_path}")
            return 0
    except Exception as e:
        print(f"❌ 오디오 길이 가져오기 중 오류: {e}")
        return 0


def create_video_from_slide_audio(slide_image, audio_file, output_video, duration):
    """슬라이드 이미지와 오디오로 비디오를 생성하는 함수"""
    import subprocess
    
    try:
        cmd = [
            'ffmpeg', '-y',  # 덮어쓰기 허용
            '-loop', '1', '-i', slide_image,  # 슬라이드 이미지
            '-i', audio_file,  # 오디오 파일
            '-c:v', 'libx264',  # 비디오 코덱
            '-t', str(duration),  # 길이
            '-pix_fmt', 'yuv420p',  # 픽셀 포맷
            '-vf', 'scale=1920:1080:force_original_aspect_ratio=increase,crop=1920:1080',  # 비디오 필터
            '-c:a', 'aac', '-b:a', '128k',  # 오디오 코덱
            output_video
        ]
        
        result = subprocess.run(cmd, capture_output=True, text=True)
        
        if result.returncode == 0:
            return True
        else:
            print(f"❌ 비디오 생성 실패: {output_video}")
            print(f"오류: {result.stderr}")
            return False
            
    except Exception as e:
        print(f"❌ 비디오 생성 중 오류: {e}")
        return False


def create_presentation_video(file_path, create_video=True):
    """발표 영상을 생성하는 함수 (Python + ffmpeg 직접 사용)"""
    if not create_video:
        return
    
    print("\n🎬 발표 영상 생성 시작...")
    
    try:
        # 필요한 디렉토리 생성
        os.makedirs("slides", exist_ok=True)
        os.makedirs("output_videos", exist_ok=True)
        
        # PDF 파일에서 페이지 이미지 추출
        print("📊 PDF에서 페이지 이미지 추출 중...")
        slide_images = extract_pages_from_pdf(file_path)
        
        if not slide_images:
            print("❌ 슬라이드 이미지 추출 실패")
            return
        
        # 스크립트 폴더 찾기
        base_name = os.path.splitext(os.path.basename(file_path))[0]
        scripts_folder = f"{base_name}_scripts"
        audio_folder = f"{scripts_folder}_audio"
        
        if not os.path.exists(scripts_folder):
            print(f"❌ 스크립트 폴더를 찾을 수 없습니다: {scripts_folder}")
            return
        
        if not os.path.exists(audio_folder):
            print(f"❌ 오디오 폴더를 찾을 수 없습니다: {audio_folder}")
            return
        
        print(f"📝 스크립트 폴더: {scripts_folder}")
        print(f"🎵 오디오 폴더: {audio_folder}")
        
        # 스크립트 파일 개수 확인
        script_files = [f for f in os.listdir(scripts_folder) if f.endswith('.txt')]
        script_count = len(script_files)
        print(f"📝 총 {script_count}개의 스크립트 파일 발견")
        
        # 각 슬라이드와 음성을 개별 비디오로 변환
        created_videos = []
        
        for i in range(1, script_count + 1):
            print(f"🎥 슬라이드 {i} 비디오 생성 중...")
            
            # 파일 경로 설정
            slide_image = f"slides/slide_{i}.png"
            audio_file = os.path.join(audio_folder, f"{base_name}_slide_{i}_script.wav")
            output_video = f"output_videos/video_{i}.mp4"
            
            # 파일 존재 확인
            if not os.path.exists(slide_image):
                print(f"❌ 슬라이드 이미지를 찾을 수 없습니다: {slide_image}")
                continue
            
            if not os.path.exists(audio_file):
                print(f"❌ 오디오 파일을 찾을 수 없습니다: {audio_file}")
                continue
            
            # 오디오 길이 가져오기
            duration = get_audio_duration(audio_file)
            if duration == 0:
                print(f"❌ 오디오 파일의 길이를 가져올 수 없습니다: {audio_file}")
                continue
            
            print(f"   오디오 길이: {duration}초")
            
            # 비디오 생성
            if create_video_from_slide_audio(slide_image, audio_file, output_video, duration):
                created_videos.append(output_video)
                print(f"✅ 슬라이드 {i} 비디오 생성 완료: {output_video}")
            else:
                print(f"❌ 슬라이드 {i} 비디오 생성 실패")
        
        # 모든 비디오를 합쳐서 최종 파일 생성
        if created_videos:
            print("🔗 모든 비디오 합치는 중...")
            
            # 비디오 파일 목록 생성
            video_list_file = "video_list.txt"
            with open(video_list_file, 'w', encoding='utf-8') as f:
                for video in created_videos:
                    f.write(f"file '{video}'\n")
            
            # 최종 비디오 파일명 생성
            final_video = f"{base_name}_발표영상.mp4"
            
            # 최종 비디오 생성
            import subprocess
            cmd = ['ffmpeg', '-y', '-f', 'concat', '-safe', '0', '-i', video_list_file, '-c', 'copy', final_video]
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode == 0:
                print(f"🎉 발표 영상 생성이 완료되었습니다: {final_video}")
                
                # 파일 크기 확인
                if os.path.exists(final_video):
                    file_size = os.path.getsize(final_video) / (1024 * 1024)  # MB
                    print(f"📁 최종 파일 크기: {file_size:.1f}MB")
            else:
                print("❌ 최종 비디오 합치기 실패")
                print(f"오류: {result.stderr}")
            
            # 임시 파일 정리
            if os.path.exists(video_list_file):
                os.remove(video_list_file)
        else:
            print("❌ 합칠 비디오 파일이 없습니다")
        
        print("🏁 영상 생성 작업 완료!")
        
    except Exception as e:
        print(f"❌ 비디오 생성 중 오류 발생: {e}")


def show_usage_examples():
    """사용 예제를 보여주는 함수"""
    print("""
🎤 Zonos 보이스 클로닝 발표 오디오 생성기

📋 기능:
- PDF 파일에서 페이지별로 간결한 발표 대본 생성
- 보이스 클로닝을 통한 개인화된 음성 생성
- 자동 발표 영상 생성
- 한국어 지원

📁 결과 파일:
- 발표자료_scripts/     # 스크립트 텍스트 파일들
- 발표자료_scripts_audio/ # 생성된 오디오 파일들
- 발표자료_발표영상.mp4  # 최종 발표 영상

🎯 사용법:
1. sample_voice.wav 파일을 준비 (10-30초 개인 음성 샘플)
2. PDF 파일을 같은 폴더에 배치
3. 스크립트 실행
""")


# --- 메인 코드 실행 ---
if __name__ == "__main__":
    # 사용 예제 보기
    show_usage_examples()
    
    # 처리할 파일 경로 (PDF)
    file_path = "중간보고서 미팅.pdf"  # PDF 파일
    
    # 간단한 프롬프트 (두 문장 내로)
    user_prompt = "핵심 내용만 간결하게 설명해주세요."
    
    # 보이스 클로닝 설정
    speaker_audio_path = "sample_voice.wav"  # 샘플 오디오 파일 경로
    language = "ko"  # 언어 설정
    generate_audio = True  # 음성 생성 여부
    
    # 통합 함수 실행 (스크립트 생성 + 음성 생성)
    generate_presentation_script(
        file_path=file_path, 
        prompt=user_prompt,
        speaker_audio_path=speaker_audio_path if os.path.exists(speaker_audio_path) else None,
        language=language,
        generate_audio=generate_audio
    )
    
    # 발표 영상 생성 (선택적)
    create_video = True  # 영상 생성을 원하지 않으면 False로 변경
    if create_video:
        create_presentation_video(file_path, create_video=True)